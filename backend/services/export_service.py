"""
Export Service - handles PPTX and PDF export
Based on demo.py create_pptx_from_images()
"""
import io
import json
import logging
import os
from typing import List

from PIL import Image
from .ai_providers.image.ppt_agent import SlideRenderer
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


class ExportService:
    """Service for exporting presentations"""
    
    @staticmethod
    def create_pptx_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PPTX file from image paths
        Based on demo.py create_pptx_from_images()
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 (width 10 inches, height 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Add each image as a slide
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout (layout 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill entire slide
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def create_pdf_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PDF file from image paths
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PDF file as bytes if output_file is None
        """
        images = []
        
        # Load all images
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            img = Image.open(image_path)
            
            # Convert to RGB if necessary (PDF requires RGB)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            images.append(img)
        
        if not images:
            raise ValueError("No valid images found for PDF export")
        
        # Save as PDF
        if output_file:
            images[0].save(
                output_file,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            return None
        else:
            # Save to bytes
            pdf_bytes = io.BytesIO()
            images[0].save(
                pdf_bytes,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            pdf_bytes.seek(0)
            return pdf_bytes.getvalue()

    """Service for exporting presentations"""

    @staticmethod
    def create_pptx_from_jsons(json_paths: List[str], output_file: str = None, assets_dir: str = "assets") -> bytes:
        """
        [推荐] 通过中间 JSON 数据重新渲染合并 PPTX
        这种方式比合并 PPTX 文件更稳定，且能保证背景和样式的一致性。

        Args:
            json_paths: JSON 文件路径列表 (包含 final_plan 和 asset_map)
            output_file: 输出路径
            assets_dir: 资源文件夹路径 (用于 SlideRenderer 加载背景图和图标)
        """
        # 1. 创建一个全新的 PPT 对象 (作为所有页面的容器)
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        logger.info(f"🚀 开始合并 {len(json_paths)} 个页面...")

        # 2. 遍历每一个 JSON 文件
        for idx, json_path in enumerate(json_paths):
            if not os.path.exists(json_path):
                logger.warning(f"❌ JSON not found: {json_path}")
                continue

            try:
                # 读取数据
                with open(json_path, 'r', encoding='utf-8') as f:
                    page_data = json.load(f)

                final_plan = page_data.get('final_plan')
                asset_map = page_data.get('asset_map', {})

                if not final_plan:
                    logger.warning(f"⚠️ Skipping {json_path}: No final_plan found.")
                    continue

                # 创建新的一页
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

                # 3. 调用渲染引擎 (SlideRenderer)
                # 注意：这里我们复用了之前的 SlideRenderer 类
                # 它会负责画背景、标题、网格、时间轴等所有元素
                renderer = SlideRenderer(prs, slide, assets_dir=assets_dir)
                renderer.dispatch(final_plan, asset_map)

                logger.info(f"✅ Page {idx + 1} rendered successfully.")

            except Exception as e:
                logger.error(f"❌ Failed to render page from {json_path}: {e}")

        # 4. 保存结果
        if output_file:
            prs.save(output_file)
            logger.info(f"🎉 Merged PPT saved to: {output_file}")
            return None
        else:
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()

    # -------------------------------------------------------------------------
    # 下面保留旧的基于图片的方法，以备不时之需
    # -------------------------------------------------------------------------
