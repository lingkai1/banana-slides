import json
import os
import math
import time
import requests
import logging
import tempfile
from typing import Optional, List, Dict, Any
from io import BytesIO
from openai import OpenAI
import httpx

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

from .base import ImageProvider

logger = logging.getLogger(__name__)

# Try to import win32com for PPT export
try:
    import win32com.client
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False

# ==========================================
# 🧠 Class 1: Planner Agent (策划)
# ==========================================
class PlannerAgent:
    def __init__(self, client: OpenAI):
        self.client = client

    def generate_plan(self, user_input: str) -> Optional[Dict[str, Any]]:
        logger.info(f"🧠 [Planner] 正在分析文本并进行多源数据融合...")
        json_schema = """{
          "meta": {"layout_type": "string", "theme": "tech_blue"},
          "content": {
            "main_title": "string", "subtitle": "string",
            "items": [{
                "id": "string", "title": "string", "desc": "string",
                "specs": {"Key": "Value"}, "tags": ["string"]
            }]
          },
          "assets": {"images": [{"target_id": "string", "prompt": "string", "local_path": null}]}
        }"""
        system_prompt = f"""你是一个高阶信息架构师。将用户输入融合为结构化数据。如果存在表格/对比数据，必须提取到 items 的 `specs` 字典中。为每个 item 生成 3D Tech Blue 风格的生图 Prompt。输出纯 JSON：{json_schema}"""
        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_input}],
                temperature=0.1
            )
            content = response.choices[0].message.content
            # Clean up markdown code blocks if present
            content = content.replace("```json", "").replace("```", "").strip()
            return json.loads(content)
        except Exception as e:
            logger.error(f"❌ 策划阶段出错: {e}")
            return None


# ==========================================
# 🏭 Class 2: Production Agent (生产)
# ==========================================
class ProductionAgent:
    def __init__(self, client: OpenAI, assets_dir: str, use_mock: bool = True):
        self.client = client
        self.use_mock = use_mock
        self.assets_dir = assets_dir
        if not os.path.exists(self.assets_dir):
            os.makedirs(self.assets_dir)

    def _create_local_mock_image(self, prompt: str, filename: str) -> str:
        filepath = os.path.join(self.assets_dir, filename)
        img = Image.new('RGB', (1024, 1024), color=(10, 30, 60))
        d = ImageDraw.Draw(img)
        d.rectangle([50, 50, 974, 974], outline=(0, 255, 255), width=5)
        d.ellipse([300, 300, 724, 724], outline=(255, 255, 255), width=2)
        # Add text to mock image to identify it
        try:
            # Try to use a default font
            d.text((320, 500), "MOCK IMAGE", fill=(255, 255, 255))
        except:
            pass

        img.save(filepath)
        logger.info(f"🎨 [Mock] 本地已生成: {filepath}")
        return filepath

    def _generate_dalle_image(self, prompt: str, filename: str) -> Optional[str]:
        logger.info(f"🎨 [DALL-E] 请求生成: {prompt[:30]}...")
        try:
            url = self.client.images.generate(model="dall-e-3", prompt=prompt, size="1024x1024", n=1).data[0].url
            res = requests.get(url)
            if res.status_code == 200:
                filepath = os.path.join(self.assets_dir, filename)
                with open(filepath, 'wb') as f: f.write(res.content)
                logger.info(f"💾 [DALL-E] 已下载: {filepath}")
                return filepath
        except Exception as e:
            logger.error(f"❌ 生图失败: {e}")
        return None

    def produce_assets(self, plan_data: Dict[str, Any]) -> Dict[str, Any]:
        logger.info(f"🏭 [Production] 开始生产图片素材...")
        for img in plan_data.get('assets', {}).get('images', []):
            fname = f"{img.get('target_id')}_{int(time.time())}.png"
            path = self._create_local_mock_image(img.get('prompt'),
                                                 fname) if self.use_mock else self._generate_dalle_image(
                img.get('prompt'), fname)
            if path: img['local_path'] = path
        return plan_data


# ==========================================
# 🔨 Class 3: Coder Agent (渲染)
# ==========================================
class SlideRenderer:
    def __init__(self, prs: Presentation, slide):
        self.slide = slide
        self.prs = prs
        self.W = prs.slide_width
        self.H = prs.slide_height
        self.C_BG = RGBColor(10, 25, 47)
        self.C_ACCENT = RGBColor(0, 255, 255)
        self.C_CARD = RGBColor(23, 42, 69)
        self.C_BORDER = RGBColor(45, 65, 95)
        self.C_TX_H = RGBColor(255, 255, 255)
        self.C_TX_B = RGBColor(170, 190, 210)

    def setup_base(self):
        bg = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.C_BG

    def draw_header(self, title, subtitle):
        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.C_ACCENT
        tb = self.slide.shapes.add_textbox(Inches(0.8), Inches(0.4), self.W - Inches(1), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.C_TX_H
        if subtitle:
            tb_s = self.slide.shapes.add_textbox(Inches(0.8), Inches(1.1), self.W - Inches(1), Inches(0.6))
            p_s = tb_s.text_frame.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(18)
            p_s.font.color.rgb = self.C_ACCENT

    def _greedy_grid(self, count, start_y):
        margin = Inches(0.5)
        gap = Inches(0.3)
        if count == 3:
            c, r = 3, 1
        elif count == 4:
            c, r = 2, 2
        elif count in [5, 6]:
            c, r = 3, 2
        else:
            c = 3; r = math.ceil(count / c)
        cw = (self.W - margin * 2 - gap * (c - 1)) / c
        ch = (self.H - start_y - margin - gap * (r - 1)) / r
        return [(margin + (i % c) * (cw + gap), start_y + (i // c) * (ch + gap), cw, ch) for i in range(count)]

    def render_grid(self, items, asset_map):
        slots = self._greedy_grid(len(items), Inches(1.8))
        for i, (x, y, w, h) in enumerate(slots):
            item = items[i]
            specs = item.get('specs', {})
            has_specs = len(specs) > 0
            card = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = self.C_CARD
            card.line.color.rgb = self.C_BORDER
            isz = Inches(0.8) if has_specs else Inches(1.2)
            cy = y + Inches(0.2)
            path = asset_map.get(item['id'])
            if path and os.path.exists(path):
                self.slide.shapes.add_picture(path, x + Inches(0.2), cy, isz, isz); tx = x + isz + Inches(0.4)
            else:
                tx = x + Inches(0.2)
            tb_t = self.slide.shapes.add_textbox(tx, cy, w - (tx - x), Inches(0.6))
            p = tb_t.text_frame.paragraphs[0]
            p.text = item.get('title', '')
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = self.C_TX_H
            tb_d = self.slide.shapes.add_textbox(tx, cy + Inches(0.4), w - (tx - x) - Inches(0.2), Inches(0.8))
            tb_d.text_frame.word_wrap = True
            p2 = tb_d.text_frame.paragraphs[0]
            p2.text = item.get('desc', '')
            p2.font.size = Pt(14)
            p2.font.color.rgb = self.C_TX_B
            if has_specs:
                sy = cy + Inches(1.2)
                sep = self.slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x + Inches(0.2), sy, w - Inches(0.4), 0)
                sep.line.color.rgb = self.C_ACCENT
                sep.line.dash_style = 1
                rh = (h - (sy - y) - Inches(0.3)) / len(specs)
                for idx, (k, v) in enumerate(specs.items()):
                    cur_y = sy + Inches(0.1) + idx * rh
                    tb_k = self.slide.shapes.add_textbox(x + Inches(0.2), cur_y, Inches(1.5), rh)
                    pk = tb_k.text_frame.paragraphs[0]
                    pk.text = f"● {k}"
                    pk.font.bold = True
                    pk.font.size = Pt(12)
                    pk.font.color.rgb = self.C_ACCENT
                    tb_v = self.slide.shapes.add_textbox(x + Inches(1.6), cur_y, w - Inches(1.8), rh)
                    tb_v.text_frame.word_wrap = True
                    tb_v.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    pv = tb_v.text_frame.paragraphs[0]
                    pv.text = str(v)
                    pv.font.size = Pt(12)
                    pv.font.color.rgb = RGBColor(200, 200, 200)

    def dispatch(self, final_plan, asset_map):
        self.setup_base()
        self.draw_header(final_plan['content'].get('main_title', ''), final_plan['content'].get('subtitle', ''))
        self.render_grid(final_plan['content'].get('items', []), asset_map)


# ==========================================
# 📸 Class 4: Exporter (导出图片)
# ==========================================
class PPTExporter:
    def __init__(self):
        self.enabled = WIN32_AVAILABLE

    def export_first_slide_as_jpg(self, pptx_path) -> Optional[str]:
        """使用 PowerPoint 后台将第一页导出为 JPG"""
        if not self.enabled:
            logger.warning("⚠️ 跳过图片导出：需要 Windows 环境并安装 pywin32。")
            return None

        # 1. 准备路径 (COM 需要绝对路径)
        abs_pptx_path = os.path.abspath(pptx_path)
        base, _ = os.path.splitext(abs_pptx_path)
        abs_jpg_path = f"{base}.jpg"

        logger.info(f"📸 [Exporter] 正在调用 PowerPoint 后台渲染图片...")
        logger.info(f"   - 源文件: {abs_pptx_path}")

        powerpoint = None
        presentation = None
        try:
            # 2. 启动 PowerPoint (隐藏模式)
            # 使用 EnsureDispatch 可以应对 PPT 已经打开的情况
            try:
                # 优先尝试动态绑定 (Dispatch)，因为它不依赖缓存，更稳定
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            except AttributeError:
                # 如果失败，再尝试强制绑定
                powerpoint = win32com.client.EnsureDispatch("PowerPoint.Application")

            # 尝试隐藏窗口
            try:
                powerpoint.Visible = False
            except:
                pass

            # 3. 打开 PPTX (只读模式，不显示窗口)
            presentation = powerpoint.Presentations.Open(abs_pptx_path, ReadOnly=True, WithWindow=False)

            # 4. 导出第一张幻灯片
            presentation.Slides(1).Export(abs_jpg_path, FilterName="JPG")
            logger.info(f"✅ 图片导出成功: {abs_jpg_path}")
            return abs_jpg_path

        except Exception as e:
            logger.error(f"❌ 图片导出失败 (COM错误): {e}")
            logger.info("提示: 请确保已安装 Microsoft PowerPoint 且未处于假死状态。")
            return None
        finally:
            # 5. 清理资源
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            if powerpoint:
                del powerpoint


class PPTAgentImageProvider(ImageProvider):
    """
    Generates an image by creating a PowerPoint slide and exporting it.
    Uses GPT-4o for planning, DALL-E 3 (or mock) for assets, and python-pptx for rendering.
    """

    def __init__(self, api_key: str, api_base: str = None, model: str = None):
        """
        Initialize PPT Agent Image Provider
        """
        self.api_key = api_key
        self.api_base = api_base
        # Ignore model parameter as this agent uses specific models internally (gpt-4o, dall-e-3)

        custom_http_client = httpx.Client(verify=False)
        self.client = OpenAI(
            api_key=self.api_key,
            base_url=self.api_base,
            http_client=custom_http_client
        )

        # Use mock images by default to avoid burning credits during dev/test
        # Ideally this should be configurable via env or param
        self.use_mock_images = os.getenv("USE_MOCK_IMAGES", "True").lower() == "true"

        # Directory for assets and temp files
        self.work_dir = os.path.join(os.getcwd(), "assets", "ppt_agent_work")
        if not os.path.exists(self.work_dir):
            os.makedirs(self.work_dir)

    def generate_image(
        self,
        prompt: str,
        ref_images: Optional[List[Image.Image]] = None,
        aspect_ratio: str = "16:9",
        resolution: str = "2K"
    ) -> Optional[Image.Image]:
        """
        Generate image by creating a PPT and exporting it.
        """
        logger.info(f"=== 🤖 PPT Agent Image Generation Start ===")
        logger.info(f"Input Prompt: {prompt[:50]}...")

        # Unique session ID for this generation
        session_id = f"ppt_{int(time.time())}"
        session_dir = os.path.join(self.work_dir, session_id)
        if not os.path.exists(session_dir):
            os.makedirs(session_dir)

        try:
            # 1. Planning
            planner = PlannerAgent(self.client)
            plan = planner.generate_plan(prompt)
            if not plan:
                logger.error("Planning failed.")
                return None

            # 2. Production
            producer = ProductionAgent(self.client, assets_dir=session_dir, use_mock=self.use_mock_images)
            final_plan = producer.produce_assets(plan)

            # 3. Coding & Rendering
            logger.info(f"🔨 [Coder] 正在构建 PPT...")
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            asset_map = {img['target_id']: img['local_path'] for img in final_plan['assets'].get('images', []) if
                         img.get('local_path')}
            renderer = SlideRenderer(prs, slide)
            renderer.dispatch(final_plan, asset_map)

            # Save PPTX
            output_ppt_name = os.path.join(session_dir, "final_output.pptx")
            prs.save(output_ppt_name)
            logger.info(f"🎉 PPT 已生成: {output_ppt_name}")

            # 4. Export
            exporter = PPTExporter()
            jpg_path = exporter.export_first_slide_as_jpg(output_ppt_name)

            if jpg_path and os.path.exists(jpg_path):
                # Load image into memory
                img = Image.open(jpg_path)
                # Force load so we can close file
                img.load()
                return img
            else:
                if not WIN32_AVAILABLE:
                    # Fallback for non-Windows environments (like Linux Sandbox)
                    logger.warning("Mocking export for Linux environment")
                    mock_img = Image.new('RGB', (1920, 1080), color=(10, 25, 47))
                    d = ImageDraw.Draw(mock_img)
                    d.text((800, 500), "PPT Generated Successfully", fill=(0, 255, 255))
                    d.text((700, 550), "(Image Export requires Windows + PowerPoint)", fill=(255, 255, 255))
                    d.text((850, 600), f"Saved at: {output_ppt_name}", fill=(170, 190, 210))
                    return mock_img

                return None

        except Exception as e:
            logger.error(f"❌ PPT Agent 流程失败: {e}", exc_info=True)
            return None
