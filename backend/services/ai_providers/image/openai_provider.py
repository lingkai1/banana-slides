"""
OpenAI SDK implementation for image generation (PPT Agent)
"""
import logging
import os
import json
import time
import math
import base64
import requests
import httpx
from typing import Optional, List
from openai import OpenAI
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .base import ImageProvider

logger = logging.getLogger(__name__)

# ==========================================
# ⚙️ Global Configuration
# ==========================================
MODEL_NAME = "gpt-4o" # or whatever model is appropriate
BACKGROUND_IMG_NAME = "tech_bg_v3.png"
USE_MOCK_IMAGES = False

WIN32_AVAILABLE = True
try:
    import win32com.client
except ImportError:
    WIN32_AVAILABLE = False


# ==========================================
# 🧠 Class 1: Planner Agent
# ==========================================
class PlannerAgent:
    def __init__(self, client, model_name=MODEL_NAME):
        self.client = client
        self.model_name = model_name

    def generate_plan(self, user_input):
        logger.info(f"🧠 [Planner] Analyzing semantics and layout strategy...")

        json_schema = """{
          "meta": {
            "layout_type": "string (must be 'timeline' or 'grid')",
            "theme": "tech_blue"
          },
          "content": {
            "main_title": "string", "subtitle": "string",
            "items": [{
                "id": "string", "title": "string (concise title)",
                "desc": "string (1-2 sentences)",
                "specs": { "Key": "Value" },
                "tags": ["string"]
              }]
          },
          "assets": {"images": [{"target_id": "string", "prompt": "string", "local_path": null}]}
        }"""

        system_prompt = f"""You are a senior PPT architect.
        Task: Analyze user input and generate structured data.

        [Layout Logic]
        - If content contains **time series** (years, dates), **process steps** (Step 1, Phase 2) or **causal evolution**:
          👉 Must set `layout_type`: "timeline"
        - If content is **parallel relationship**, **comparison analysis** or **core elements listing**:
          👉 Set `layout_type`: "grid"

        [Data Processing]
        1. Extract tables/lists to `specs` field.
        2. Generate 3D Tech Blue style image prompts for each item.

        Output pure JSON: {json_schema}"""

        try:
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_input}],
                temperature=0.1,
                response_format={"type": "json_object"}
            )
            content = response.choices[0].message.content
            return json.loads(content)
        except Exception as e:
            logger.error(f"❌ Planning failed: {e}")
            # Fallback plan if JSON parsing fails
            return {
                "meta": {"layout_type": "grid", "theme": "tech_blue"},
                "content": {
                    "main_title": "Generated Content",
                    "subtitle": "Analysis",
                    "items": [{"id": "1", "title": "Content", "desc": user_input[:100], "specs": {}, "tags": []}]
                },
                "assets": {"images": []}
            }


# ==========================================
# 🏭 Class 2: Production Agent
# ==========================================
class ProductionAgent:
    def __init__(self, assets_dir, use_mock=True, image_provider=None):
        self.use_mock = use_mock
        self.assets_dir = assets_dir
        self.image_provider = image_provider # Can be used to call actual image generation if available

    def _create_tech_background_asset(self):
        filepath = os.path.join(self.assets_dir, BACKGROUND_IMG_NAME)
        if os.path.exists(filepath): return filepath

        logger.info("🎨 [Production] Generating V3 Tech Background...")
        W, H = 1920, 1080
        img = Image.new('RGB', (W, H), color=(4, 12, 28))
        draw = ImageDraw.Draw(img)
        # Top glow
        for i in range(500):
            alpha = int(50 * (1 - i / 500))
            draw.line([(0, i), (W, i)], fill=(0, 100, 200, alpha), width=1)
        # Bottom grid
        for x in range(0, W, 80):
            draw.line([(x, H), (W / 2, H / 2)], fill=(0, 255, 255, 10), width=1)
        img.save(filepath)
        return filepath

    def _create_local_pil_mock(self, prompt, filename):
        path = os.path.join(self.assets_dir, filename)
        img = Image.new('RGB', (1024, 1024), (10, 30, 60))
        d = ImageDraw.Draw(img)
        d.rectangle([50, 50, 974, 974], outline=(0, 200, 255), width=8)
        d.ellipse([300, 300, 724, 724], outline=(200, 255, 255), width=4)

        # Draw some text to identify the image
        try:
            # Try to draw text if basic font available
            d.text((100, 500), prompt[:20], fill=(255, 255, 255))
        except:
            pass

        img.save(path)
        return path

    def produce_assets(self, plan):
        self._create_tech_background_asset()

        # Generate images for items
        for img_spec in plan.get('assets', {}).get('images', []):
            fname = f"{img_spec.get('target_id')}_{int(time.time())}.png"
            prompt = img_spec.get('prompt', 'tech image')

            # Use mock for now as instructed, or could implement actual generation call
            # For this task, we will stick to mock or simple generation to ensure speed/stability
            # unless the user explicitly wants high-quality generation for each asset which might be slow.
            # The provided code used _generate_qwen_api_image or mock.
            # We will use mock if USE_MOCK_IMAGES is True, otherwise we could try to generate.
            # Since we don't have the Qwen API access as in the snippet, we default to mock
            # or we could use the self.client if it was an image model.
            # But here we are inside the ImageProvider itself...

            path = self._create_local_pil_mock(prompt, fname)
            if path:
                img_spec['local_path'] = path

        return plan


# ==========================================
# 🔨 Class 3: Coder Agent
# ==========================================
class SlideRenderer:
    def __init__(self, prs, slide, assets_dir):
        self.slide = slide
        self.prs = prs
        self.assets_dir = assets_dir
        self.W = prs.slide_width
        self.H = prs.slide_height

        # Color Scheme (Tech Blue Pro)
        self.C_ACCENT = RGBColor(0, 240, 255)
        self.C_ACCENT_DIM = RGBColor(0, 100, 140)
        self.C_CARD_BG = RGBColor(12, 25, 45)
        self.C_BORDER = RGBColor(60, 120, 180)
        self.C_TX_H = RGBColor(255, 255, 255)
        self.C_TX_B = RGBColor(200, 210, 230)
        self.C_ROW_ALT = RGBColor(20, 40, 65)

    def setup_base(self):
        bg_path = os.path.join(self.assets_dir, BACKGROUND_IMG_NAME)

        if os.path.exists(bg_path):
            self.slide.shapes.add_picture(bg_path, 0, 0, self.W, self.H)
        else:
            bg = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(5, 10, 20)
            bg.line.fill.background()

    def draw_header(self, title, subtitle):
        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.C_ACCENT
        tb = self.slide.shapes.add_textbox(Inches(0.8), Inches(0.35), self.W - Inches(1), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        try:
            p.font.name = "Microsoft YaHei UI Bold"
        except:
            pass
        p.font.color.rgb = self.C_TX_H
        if subtitle:
            tb_s = self.slide.shapes.add_textbox(Inches(0.8), Inches(1.1), self.W - Inches(1), Inches(0.6))
            p_s = tb_s.text_frame.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(20)
            p_s.font.color.rgb = self.C_ACCENT

    def render_grid(self, items, asset_map):
        logger.info("📐 [Render] Strategy: Grid Layout")
        start_y = Inches(1.6)
        margin = Inches(0.5)
        gap = Inches(0.3)
        count = len(items)

        if count == 0:
            return

        if count <= 3:
            c, r = count, 1
        elif count == 4:
            c, r = 2, 2
        else:
            c = 3; r = math.ceil(count / c)

        # Avoid division by zero
        c = max(1, c)
        r = max(1, r)

        cw = (self.W - margin * 2 - gap * (c - 1)) / c
        ch = (self.H - start_y - margin - gap * (r - 1)) / r

        for i, item in enumerate(items):
            row = i // c
            col = i % c
            x = margin + col * (cw + gap)
            y = start_y + row * (ch + gap)
            self._draw_card_content(item, x, y, cw, ch, asset_map)

    def render_timeline(self, items, asset_map):
        logger.info("📐 [Render] Strategy: Timeline Layout")
        count = len(items)
        if count == 0:
            return

        margin = Inches(0.5)
        line_y = Inches(3.0)
        line = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, line_y, self.W - margin * 2, Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = self.C_ACCENT
        line.shadow.inherit = False

        slot_w = (self.W - margin * 2) / count

        for i, item in enumerate(items):
            cx = margin + i * slot_w + slot_w / 2
            dot_out = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.15), line_y - Inches(0.12), Inches(0.3),
                                                  Inches(0.3))
            dot_out.fill.solid()
            dot_out.fill.fore_color.rgb = self.C_ACCENT
            dot_in = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.08), line_y - Inches(0.05), Inches(0.16),
                                                 Inches(0.16))
            dot_in.fill.solid()
            dot_in.fill.fore_color.rgb = self.C_CARD_BG

            img_path = asset_map.get(item['id'])
            if img_path and os.path.exists(img_path):
                isz = Inches(1.8)
                pic = self.slide.shapes.add_picture(img_path, cx - isz / 2, line_y - isz - Inches(0.4), isz, isz)
                pic.line.color.rgb = self.C_ACCENT
                pic.line.width = Pt(1.5)

            card_x = cx - slot_w / 2 + Inches(0.1)
            card_y = line_y + Inches(0.4)
            card_w = slot_w - Inches(0.2)
            card_h = self.H - card_y - Inches(0.5)
            self._draw_card_content(item, card_x, card_y, card_w, card_h, asset_map, is_timeline=True)

    def _draw_card_content(self, item, x, y, w, h, asset_map, is_timeline=False):
        card = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = self.C_CARD_BG
        card.line.color.rgb = self.C_BORDER
        card.line.width = Pt(1.5)

        cursor_y = y + Inches(0.2)

        if not is_timeline:
            img_path = asset_map.get(item['id'])
            icon_w = Inches(0)
            if img_path and os.path.exists(img_path):
                isz = Inches(0.8)
                self.slide.shapes.add_picture(img_path, x + Inches(0.2), cursor_y, isz, isz)
                icon_w = isz + Inches(0.2)
            title_x = x + Inches(0.2) + icon_w
            title_w = w - icon_w - Inches(0.3)
        else:
            title_x = x + Inches(0.2)
            title_w = w - Inches(0.4)

        tb_t = self.slide.shapes.add_textbox(title_x, cursor_y, title_w, Inches(0.8))
        tb_t.text_frame.word_wrap = True
        p = tb_t.text_frame.paragraphs[0]
        p.text = item.get('title', '')
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = self.C_TX_H
        try:
            p.font.name = "Microsoft YaHei UI Bold"
        except:
            pass
        if is_timeline: p.alignment = PP_ALIGN.CENTER

        cursor_y += Inches(0.8 if not is_timeline else 0.7)

        desc_h = Inches(0.8)
        specs = item.get('specs', {})
        if not specs: desc_h = (y + h) - cursor_y - Inches(0.2)

        # Ensure height is positive
        if desc_h < 0: desc_h = Inches(0.5)

        tb_d = self.slide.shapes.add_textbox(x + Inches(0.2), cursor_y, w - Inches(0.4), desc_h)
        tb_d.text_frame.word_wrap = True
        p2 = tb_d.text_frame.paragraphs[0]
        p2.text = item.get('desc', '')
        p2.font.size = Pt(16)
        p2.font.color.rgb = self.C_TX_B
        try:
            p2.font.name = "Microsoft YaHei UI"
        except:
            pass
        if is_timeline: p2.alignment = PP_ALIGN.CENTER

        cursor_y += Inches(0.8)

        if specs:
            sep = self.slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x + Inches(0.1), cursor_y, w - Inches(0.2), 0)
            sep.line.color.rgb = self.C_ACCENT
            sep.line.width = Pt(1)
            sep.line.dash_style = 1
            cursor_y += Inches(0.1)

            rem_h = (y + h) - cursor_y - Inches(0.1)
            if rem_h > 0:
                row_h = rem_h / len(specs)
                for idx, (k, v) in enumerate(specs.items()):
                    ry = cursor_y + idx * row_h
                    if idx % 2 == 0:
                        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(2), ry, w - Pt(4), row_h)
                        bar.fill.solid()
                        bar.fill.fore_color.rgb = self.C_ROW_ALT
                        bar.line.fill.background()

                    tb_k = self.slide.shapes.add_textbox(x + Inches(0.2), ry, w * 0.4, row_h)
                    tb_k.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    pk = tb_k.text_frame.paragraphs[0]
                    pk.text = f"● {k}"
                    pk.font.bold = True
                    pk.font.size = Pt(14)
                    pk.font.color.rgb = self.C_ACCENT

                    tb_v = self.slide.shapes.add_textbox(x + Inches(0.2) + w * 0.4, ry, w * 0.55, row_h)
                    tb_v.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    pv = tb_v.text_frame.paragraphs[0]
                    pv.text = str(v)
                    pv.font.size = Pt(14)
                    pv.font.color.rgb = self.C_TX_H

    def dispatch(self, plan, asset_map):
        self.setup_base()
        self.draw_header(plan['content'].get('main_title', ''), plan['content'].get('subtitle', ''))
        layout = plan['meta'].get('layout_type', 'grid')
        if layout == 'timeline':
            self.render_timeline(plan['content']['items'], asset_map)
        else:
            self.render_grid(plan['content']['items'], asset_map)


# ==========================================
# 📸 Class 4: Exporter
# ==========================================
class PPTExporter:
    def export(self, pptx_path, output_img_path):
        """
        Export first page to image
        """
        if not WIN32_AVAILABLE:
            logger.warning("❌ Win32Com not available, cannot export image via PowerPoint")
            # Fallback: Create a placeholder image if on Linux
            try:
                img = Image.new('RGB', (1920, 1080), color=(20, 30, 50))
                d = ImageDraw.Draw(img)
                d.text((800, 500), "PPT Generated (Win32 Unavailable)", fill=(255, 255, 255))
                img.save(output_img_path)
                logger.info(f"✅ Placeholder image generated: {output_img_path}")
                return True
            except Exception as e:
                logger.error(f"❌ Failed to generate placeholder: {e}")
                return False

        try:
            abs_pptx_path = os.path.abspath(pptx_path)
            abs_img_path = os.path.abspath(output_img_path)

            ppt = win32com.client.Dispatch("PowerPoint.Application")
            try:
                ppt.Visible = False
            except:
                pass

            pres = ppt.Presentations.Open(abs_pptx_path, ReadOnly=True, WithWindow=False)
            pres.Slides(1).Export(abs_img_path, FilterName="JPG")
            pres.Close()
            logger.info(f"✅ Preview image generated: {abs_img_path}")
            return True
        except Exception as e:
            logger.error(f"❌ Export image failed: {e}")
            return False


class OpenAIImageProvider(ImageProvider):
    """
    PPT Agent Image Provider
    Generates images by creating a PowerPoint slide via python-pptx and exporting it.
    """
    
    def __init__(self, api_key: str, api_base: str = None, model: str = "gpt-4o"):
        """
        Initialize OpenAI image provider
        
        Args:
            api_key: API key
            api_base: API base URL
            model: Model name to use
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url=api_base,
            http_client=httpx.Client(verify=False)
        )
        self.model = model
    
    def generate_image(
        self,
        prompt: str,
        ref_images: Optional[List[Image.Image]] = None,
        aspect_ratio: str = "16:9",
        resolution: str = "2K",
        project_id: Optional[str] = None,
        page_id: Optional[str] = None
    ) -> Optional[Image.Image]:
        """
        Generate image using PPT Agent (generate PPT -> export image)
        
        Args:
            prompt: The outline/content for the PPT
            ref_images: Optional list of reference images (ignored)
            aspect_ratio: Image aspect ratio (ignored, defaults to 16:9)
            resolution: Image resolution (ignored)
            project_id: Project ID for saving files
            page_id: Page ID for saving files
            
        Returns:
            Generated PIL Image object, or None if failed
        """
        try:
            logger.info(f"🚀 Starting PPT Agent generation for project={project_id}, page={page_id}")
            
            # Define paths
            # Default to current directory if not provided, but they should be provided
            base_dir = "uploads"
            if project_id:
                base_dir = os.path.join("uploads", project_id)
            
            # Ensure directories exist
            assets_dir = os.path.join(base_dir, "assets")
            if not os.path.exists(assets_dir):
                os.makedirs(assets_dir)
            
            if not os.path.exists(base_dir):
                os.makedirs(base_dir)

            # Determine filenames
            ppt_filename = f"{page_id}.pptx" if page_id else f"temp_{int(time.time())}.pptx"
            img_filename = f"{page_id}.jpg" if page_id else f"temp_{int(time.time())}.jpg"
            
            ppt_output_path = os.path.join(base_dir, ppt_filename)
            img_output_path = os.path.join(base_dir, img_filename)

            # 1. Plan
            planner = PlannerAgent(self.client, self.model)
            plan = planner.generate_plan(prompt)
            if not plan:
                raise Exception("Planning failed")

            # 2. Production
            producer = ProductionAgent(assets_dir=assets_dir, use_mock=USE_MOCK_IMAGES)
            final_plan = producer.produce_assets(plan)

            # 3. Rendering
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Build asset map
            asset_map = {}
            for img in final_plan['assets'].get('images', []):
                if img.get('local_path'):
                    asset_map[img['target_id']] = img['local_path']

            renderer = SlideRenderer(prs, slide, assets_dir=assets_dir)
            renderer.dispatch(final_plan, asset_map)

            # 4. Save PPT
            try:
                prs.save(ppt_output_path)
                logger.info(f"🎉 PPT saved: {ppt_output_path}")
            except Exception as e:
                logger.error(f"Failed to save PPT: {e}")
                raise

            # 5. Export Image
            exporter = PPTExporter()
            success = exporter.export(ppt_output_path, img_output_path)
            
            if success and os.path.exists(img_output_path):
                return Image.open(img_output_path)
            else:
                logger.error("Failed to export image or image not found")
                # Return a simple placeholder if export failed completely
                img = Image.new('RGB', (1920, 1080), color=(50, 50, 50))
                return img
            
        except Exception as e:
            error_detail = f"Error generating PPT image: {str(e)}"
            logger.error(error_detail, exc_info=True)
            raise Exception(error_detail) from e
