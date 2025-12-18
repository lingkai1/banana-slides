import json
import os
import math
import time
import base64
import requests
import httpx
import logging
from openai import OpenAI
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ==========================================
# ⚙️ Configuration
# ==========================================
BACKGROUND_IMG_NAME = "tech_bg_v3.png"
# Default to using mock images if we can't reach the internal API
USE_MOCK_IMAGES = os.environ.get("PPT_AGENT_USE_MOCK", "True").lower() == "true"

WIN32_AVAILABLE = True
try:
    import win32com.client
except ImportError:
    WIN32_AVAILABLE = False


# ==========================================
# 🧠 Class 1: Planner Agent
# ==========================================
class PlannerAgent:
    def __init__(self, client, model_name):
        self.client = client
        self.model_name = model_name

    def generate_plan(self, user_input):
        logger.info(f"🧠 [Planner] 正在分析语义与布局策略...")

        json_schema = """{
          "meta": {
            "layout_type": "string (只能填 'timeline' 或 'grid')",
            "theme": "tech_blue"
          },
          "content": {
            "main_title": "string", "subtitle": "string",
            "items": [{
                "id": "string", "title": "string (简练标题)",
                "desc": "string (1-2句话)",
                "specs": { "Key": "Value" },
                "tags": ["string"]
              }]
          },
          "assets": {"images": [{"target_id": "string", "prompt": "string", "local_path": null}]}
        }"""

        system_prompt = f"""你是一个高级PPT架构师。
        任务：分析用户输入，生成结构化数据。

        【布局决策逻辑】
        - 如果内容包含**时间序列**（年份、日期）、**步骤流程**（Step 1, Phase 2）或**因果演进**：
          👉 必须设置 `layout_type`: "timeline"
        - 如果内容是**并列关系**、**对比分析**或**核心要素列举**：
          👉 设置 `layout_type`: "grid"

        【数据处理】
        1. 将表格/列表提取到 `specs` 字段。
        2. 为每个 item 生成 3D Tech Blue 风格的生图 Prompt。

        输出纯 JSON：{json_schema}"""

        try:
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_input}],
                temperature=0.1
            )
            content = response.choices[0].message.content.replace("```json", "").replace("```", "").strip()
            return json.loads(content)
        except Exception as e:
            logger.error(f"❌ 策划失败: {e}")
            return None


# ==========================================
# 🏭 Class 2: Production Agent (改造：接收 assets_dir)
# ==========================================
class ProductionAgent:
    def __init__(self, assets_dir, use_mock=True, image_generator=None):
        self.use_mock = use_mock
        self.assets_dir = assets_dir  # 接收外部传入的路径
        self.image_generator = image_generator # Optional callback for generating images

    def _create_tech_background_asset(self):
        filepath = os.path.join(self.assets_dir, BACKGROUND_IMG_NAME)
        if os.path.exists(filepath): return filepath

        logger.info("🎨 [Production] 生成 V3 科技背景...")
        W, H = 1920, 1080
        img = Image.new('RGB', (W, H), color=(4, 12, 28))
        draw = ImageDraw.Draw(img)
        # 顶部光辉
        for i in range(500):
            alpha = int(50 * (1 - i / 500))
            draw.line([(0, i), (W, i)], fill=(0, 100, 200, alpha), width=1)
        # 底部网格
        for x in range(0, W, 80):
            draw.line([(x, H), (W / 2, H / 2)], fill=(0, 255, 255, 10), width=1)
        img.save(filepath)
        return filepath

    def _generate_qwen_api_image(self, prompt, filename):
        # Allow override via image_generator callback
        if self.image_generator:
            try:
                img = self.image_generator(prompt)
                if img:
                    path = os.path.join(self.assets_dir, filename)
                    img.save(path)
                    return path
            except Exception as e:
                logger.error(f"Image generator callback failed: {e}")

        # Fallback to hardcoded API or mock
        url = "http://10.155.71.211:18888/qwen_image"
        try:
            res = requests.post(url, json={"prompt": prompt}, headers={'Content-Type': 'application/json'}, timeout=10) # Reduced timeout
            if res.status_code == 200 and res.json().get("status") == "success":
                b64 = res.json().get("image_base64", "").split(",")[-1]
                path = os.path.join(self.assets_dir, filename)
                with open(path, 'wb') as f: f.write(base64.b64decode(b64))
                return path
        except Exception as e:
            logger.warning(f"Qwen API failed: {e}")

        # If API fails, fallback to mock
        logger.warning(f"Falling back to mock image for {filename}")
        return self._create_local_pil_mock(prompt, filename)

    def _create_local_pil_mock(self, prompt, filename):
        path = os.path.join(self.assets_dir, filename)
        img = Image.new('RGB', (1024, 1024), (10, 30, 60))
        d = ImageDraw.Draw(img)
        d.rectangle([50, 50, 974, 974], outline=(0, 200, 255), width=8)
        d.ellipse([300, 300, 724, 724], outline=(200, 255, 255), width=4)

        # Add simple text
        try:
            # Try to load a font, or use default
            d.text((100, 100), "Mock Image", fill=(255, 255, 255))
            d.text((100, 150), prompt[:20], fill=(255, 255, 255))
        except:
            pass

        img.save(path)
        return path

    def produce_assets(self, plan):
        self._create_tech_background_asset()
        for img in plan.get('assets', {}).get('images', []):
            fname = f"{img.get('target_id')}_{int(time.time())}.png"
            if self.use_mock:
                path = self._create_local_pil_mock(img.get('prompt'), fname)
            else:
                path = self._generate_qwen_api_image(img.get('prompt'), fname)

            if path: img['local_path'] = path
        return plan


# ==========================================
# 🔨 Class 3: Coder Agent (改造：接收 assets_dir)
# ==========================================
class SlideRenderer:
    def __init__(self, prs, slide, assets_dir):
        self.slide = slide
        self.prs = prs
        self.assets_dir = assets_dir  # 接收外部传入的路径
        self.W = prs.slide_width
        self.H = prs.slide_height

        # 配色系统 (Tech Blue Pro)
        self.C_ACCENT = RGBColor(0, 240, 255)
        self.C_ACCENT_DIM = RGBColor(0, 100, 140)
        self.C_CARD_BG = RGBColor(12, 25, 45)
        self.C_BORDER = RGBColor(60, 120, 180)
        self.C_TX_H = RGBColor(255, 255, 255)
        self.C_TX_B = RGBColor(200, 210, 230)
        self.C_ROW_ALT = RGBColor(20, 40, 65)

    def setup_base(self):
        bg_path = os.path.join(self.assets_dir, BACKGROUND_IMG_NAME)

        if os.path.exists(bg_path):
            self.slide.shapes.add_picture(bg_path, 0, 0, self.W, self.H)
        else:
            bg = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(5, 10, 20)
            bg.line.fill.background()

    def draw_header(self, title, subtitle):
        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.C_ACCENT
        tb = self.slide.shapes.add_textbox(Inches(0.8), Inches(0.35), self.W - Inches(1), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        # p.font.name = "Microsoft YaHei UI Bold" # Font might not exist on Linux
        p.font.color.rgb = self.C_TX_H
        if subtitle:
            tb_s = self.slide.shapes.add_textbox(Inches(0.8), Inches(1.1), self.W - Inches(1), Inches(0.6))
            p_s = tb_s.text_frame.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(20)
            p_s.font.color.rgb = self.C_ACCENT

    def render_grid(self, items, asset_map):
        logger.info("📐 [Render] 执行策略: Grid Layout")
        start_y = Inches(1.6)
        margin = Inches(0.5)
        gap = Inches(0.3)
        count = len(items)

        if count <= 3:
            c, r = count, 1
        elif count == 4:
            c, r = 2, 2
        else:
            c = 3; r = math.ceil(count / c)
        cw = (self.W - margin * 2 - gap * (c - 1)) / c
        ch = (self.H - start_y - margin - gap * (r - 1)) / r

        for i, item in enumerate(items):
            row = i // c;
            col = i % c
            x = margin + col * (cw + gap)
            y = start_y + row * (ch + gap)
            self._draw_card_content(item, x, y, cw, ch, asset_map)

    def render_timeline(self, items, asset_map):
        logger.info("📐 [Render] 执行策略: Timeline Layout")
        count = len(items)
        margin = Inches(0.5)
        line_y = Inches(3.0)
        line = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, line_y, self.W - margin * 2, Inches(0.06))
        line.fill.solid();
        line.fill.fore_color.rgb = self.C_ACCENT
        line.shadow.inherit = False

        slot_w = (self.W - margin * 2) / count

        for i, item in enumerate(items):
            cx = margin + i * slot_w + slot_w / 2
            dot_out = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.15), line_y - Inches(0.12), Inches(0.3),
                                                  Inches(0.3))
            dot_out.fill.solid();
            dot_out.fill.fore_color.rgb = self.C_ACCENT
            dot_in = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.08), line_y - Inches(0.05), Inches(0.16),
                                                 Inches(0.16))
            dot_in.fill.solid();
            dot_in.fill.fore_color.rgb = self.C_CARD_BG

            img_path = asset_map.get(item['id'])
            if img_path and os.path.exists(img_path):
                isz = Inches(1.8)
                pic = self.slide.shapes.add_picture(img_path, cx - isz / 2, line_y - isz - Inches(0.4), isz, isz)
                pic.line.color.rgb = self.C_ACCENT;
                pic.line.width = Pt(1.5)

            card_x = cx - slot_w / 2 + Inches(0.1)
            card_y = line_y + Inches(0.4)
            card_w = slot_w - Inches(0.2)
            card_h = self.H - card_y - Inches(0.5)
            self._draw_card_content(item, card_x, card_y, card_w, card_h, asset_map, is_timeline=True)

    def _draw_card_content(self, item, x, y, w, h, asset_map, is_timeline=False):
        card = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid();
        card.fill.fore_color.rgb = self.C_CARD_BG
        card.line.color.rgb = self.C_BORDER;
        card.line.width = Pt(1.5)

        cursor_y = y + Inches(0.2)

        if not is_timeline:
            img_path = asset_map.get(item['id'])
            icon_w = Inches(0)
            if img_path and os.path.exists(img_path):
                isz = Inches(0.8)
                self.slide.shapes.add_picture(img_path, x + Inches(0.2), cursor_y, isz, isz)
                icon_w = isz + Inches(0.2)
            title_x = x + Inches(0.2) + icon_w
            title_w = w - icon_w - Inches(0.3)
        else:
            title_x = x + Inches(0.2)
            title_w = w - Inches(0.4)

        tb_t = self.slide.shapes.add_textbox(title_x, cursor_y, title_w, Inches(0.8))
        tb_t.text_frame.word_wrap = True
        p = tb_t.text_frame.paragraphs[0]
        p.text = item.get('title', '')
        p.font.bold = True;
        p.font.size = Pt(24);
        p.font.color.rgb = self.C_TX_H
        # p.font.name = "Microsoft YaHei UI Bold"
        if is_timeline: p.alignment = PP_ALIGN.CENTER

        cursor_y += Inches(0.8 if not is_timeline else 0.7)

        desc_h = Inches(0.8)
        specs = item.get('specs', {})
        if not specs: desc_h = (y + h) - cursor_y - Inches(0.2)

        tb_d = self.slide.shapes.add_textbox(x + Inches(0.2), cursor_y, w - Inches(0.4), desc_h)
        tb_d.text_frame.word_wrap = True
        p2 = tb_d.text_frame.paragraphs[0]
        p2.text = item.get('desc', '')
        p2.font.size = Pt(16);
        p2.font.color.rgb = self.C_TX_B
        # p2.font.name = "Microsoft YaHei UI"
        if is_timeline: p2.alignment = PP_ALIGN.CENTER

        cursor_y += Inches(0.8)

        if specs:
            sep = self.slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x + Inches(0.1), cursor_y, w - Inches(0.2), 0)
            sep.line.color.rgb = self.C_ACCENT;
            sep.line.width = Pt(1);
            sep.line.dash_style = 1
            cursor_y += Inches(0.1)

            rem_h = (y + h) - cursor_y - Inches(0.1)
            if rem_h > 0:
                row_h = rem_h / len(specs)
                for idx, (k, v) in enumerate(specs.items()):
                    ry = cursor_y + idx * row_h
                    if idx % 2 == 0:
                        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(2), ry, w - Pt(4), row_h)
                        bar.fill.solid();
                        bar.fill.fore_color.rgb = self.C_ROW_ALT;
                        bar.line.fill.background()

                    tb_k = self.slide.shapes.add_textbox(x + Inches(0.2), ry, w * 0.4, row_h)
                    tb_k.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    pk = tb_k.text_frame.paragraphs[0]
                    pk.text = f"● {k}";
                    pk.font.bold = True;
                    pk.font.size = Pt(14);
                    pk.font.color.rgb = self.C_ACCENT

                    tb_v = self.slide.shapes.add_textbox(x + Inches(0.2) + w * 0.4, ry, w * 0.55, row_h)
                    tb_v.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    pv = tb_v.text_frame.paragraphs[0]
                    pv.text = str(v);
                    pv.font.size = Pt(14);
                    pv.font.color.rgb = self.C_TX_H

    def dispatch(self, plan, asset_map):
        self.setup_base()
        self.draw_header(plan['content'].get('main_title', ''), plan['content'].get('subtitle', ''))
        layout = plan['meta'].get('layout_type', 'grid')
        if layout == 'timeline':
            self.render_timeline(plan['content']['items'], asset_map)
        else:
            self.render_grid(plan['content']['items'], asset_map)


# ==========================================
# 📸 Class 4: Exporter (改造：支持自定义输出图片路径)
# ==========================================
class PPTExporter:
    def export(self, pptx_path, output_img_path):
        """
        导出第一页为图片
        :param pptx_path: PPT 文件路径
        :param output_img_path: 目标图片保存路径 (必须是绝对路径或相对路径)
        """
        if not WIN32_AVAILABLE:
            logger.warning("❌ Win32Com 不可用，正在生成占位图作为替代")
            self._create_placeholder_image(output_img_path)
            return

        try:
            # 转换为绝对路径，因为 COM 接口通常需要绝对路径
            abs_pptx_path = os.path.abspath(pptx_path)
            abs_img_path = os.path.abspath(output_img_path)

            ppt = win32com.client.Dispatch("PowerPoint.Application")
            try:
                ppt.Visible = False
            except:
                pass

            # 以只读方式打开，不显示窗口
            pres = ppt.Presentations.Open(abs_pptx_path, ReadOnly=True, WithWindow=False)

            # Export 方法导出第一张幻灯片
            # FilterName 指定格式
            pres.Slides(1).Export(abs_img_path, FilterName="JPG")

            pres.Close()
            # ppt.Quit() # 生产环境建议不频繁杀进程，或者手动管理
            logger.info(f"✅ 预览图已生成: {abs_img_path}")
        except Exception as e:
            logger.error(f"❌ 导出图片出错: {e}, 使用占位图")
            self._create_placeholder_image(output_img_path)

    def _create_placeholder_image(self, output_path):
        """生成一张占位图片，当 PowerPoint 导出不可用时使用"""
        try:
            W, H = 1280, 720
            img = Image.new('RGB', (W, H), color=(20, 30, 50))
            draw = ImageDraw.Draw(img)

            # Draw a box
            draw.rectangle([100, 100, W-100, H-100], outline=(0, 240, 255), width=5)

            # Draw text
            text = "PPT Preview Unavailable (Linux/No Win32)"
            # Try to center text
            bbox = draw.textbbox((0, 0), text)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]
            draw.text(((W-tw)/2, (H-th)/2), text, fill=(255, 255, 255))

            img.save(output_path)
            logger.info(f"Generated placeholder image at {output_path}")
        except Exception as e:
            logger.error(f"Failed to generate placeholder image: {e}")


# ==========================================
# 🚀 API Function: 对外提供的集成函数
# ==========================================
def generate_single_page_ppt(
    outline: str,
    ppt_output_path: str,
    img_output_path: str,
    assets_output_dir: str,
    client: OpenAI,
    model_name: str,
    image_generator=None
):
    """
    生成单页 PPT 的主函数
    :param outline: PPT 大纲/内容文本
    :param ppt_output_path: 输出 PPT 的路径 (e.g., "./output/result.pptx")
    :param img_output_path: 输出图片的路径 (e.g., "./output/preview.jpg")
    :param assets_output_dir: 中间素材(图片)存放目录 (e.g., "./assets")
    :param client: OpenAI client instance
    :param model_name: Model name
    :param image_generator: Optional callback to generate images (prompt -> PIL.Image)
    :return: 字典 result {"status": "success/error", "ppt_path": ..., "img_path": ...}
    """

    # 1. 确保目录存在
    if not os.path.exists(assets_output_dir):
        os.makedirs(assets_output_dir)

    ppt_dir = os.path.dirname(ppt_output_path)
    if ppt_dir and not os.path.exists(ppt_dir):
        os.makedirs(ppt_dir)

    img_dir = os.path.dirname(img_output_path)
    if img_dir and not os.path.exists(img_dir):
        os.makedirs(img_dir)

    logger.info(f"=== 🚀 开始生成 PPT Agent 任务 ===")

    # 2. Step 1: Planning
    planner = PlannerAgent(client, model_name)
    plan = planner.generate_plan(outline)
    if not plan:
        return {"status": "error", "message": "Planning failed"}

    logger.info(f"🎯 [Strategy] LLM 选定布局: {plan['meta'].get('layout_type', 'unknown').upper()}")

    # 3. Step 2: Production
    # 将 assets_output_dir 传入 ProductionAgent
    producer = ProductionAgent(
        assets_dir=assets_output_dir,
        use_mock=USE_MOCK_IMAGES,
        image_generator=image_generator
    )
    final_plan = producer.produce_assets(plan)

    # 4. Step 3: Rendering
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 构建资源映射 (使用绝对路径或相对 assets 目录的路径)
    asset_map = {}
    for img in final_plan['assets'].get('images', []):
        if img.get('local_path'):
            asset_map[img['target_id']] = img['local_path']

    # 将 assets_output_dir 传入 SlideRenderer (用于寻找背景图)
    renderer = SlideRenderer(prs, slide, assets_dir=assets_output_dir)
    renderer.dispatch(final_plan, asset_map)

    # 5. Save PPT
    try:
        prs.save(ppt_output_path)
        logger.info(f"🎉 PPT 已保存: {os.path.abspath(ppt_output_path)}")
    except Exception as e:
        return {"status": "error", "message": f"Save PPT failed: {str(e)}"}

    # 6. Step 4: Export Image
    exporter = PPTExporter()
    exporter.export(ppt_output_path, img_output_path)

    return {
        "status": "success",
        "ppt_path": os.path.abspath(ppt_output_path),
        "img_path": os.path.abspath(img_output_path)
    }
